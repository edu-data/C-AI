"""
다중 형식 리포트 생성기 - GAIM Lab v3.0
PDF, Excel, PowerPoint 형식 지원
"""

import json
from pathlib import Path
from datetime import datetime
from typing import List, Dict


class ReportGenerator:
    """다양한 형식의 리포트 생성"""
    
    def __init__(self, results_data: List[Dict], output_dir: str):
        """
        리포트 생성 초기화
        
        Args:
            results_data: 분석 결과 데이터
            output_dir: 출력 디렉토리
        """
        self.results = results_data
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
    
    def generate_excel_report(self, filename: str = "analysis_report.xlsx") -> str:
        """
        Excel 리포트 생성
        
        Args:
            filename: 출력 파일명
            
        Returns:
            생성된 파일 경로
        """
        try:
            import openpyxl
            from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
            from openpyxl.utils import get_column_letter
        except ImportError:
            return "❌ openpyxl 패키지가 필요합니다: pip install openpyxl"
        
        # 워크북 생성
        wb = openpyxl.Workbook()
        
        # 1. 요약 시트
        ws = wb.active
        ws.title = "요약"
        self._create_summary_sheet(ws)
        
        # 2. 상세 분석 시트
        ws_detail = wb.create_sheet("상세분석")
        self._create_detail_sheet(ws_detail)
        
        # 3. 교사별 분석 시트
        ws_teacher = wb.create_sheet("교사별분석")
        self._create_teacher_sheet(ws_teacher)
        
        # 4. 과목별 분석 시트
        ws_subject = wb.create_sheet("과목별분석")
        self._create_subject_sheet(ws_subject)
        
        # 파일 저장
        output_path = self.output_dir / filename
        wb.save(output_path)
        
        return f"✅ Excel 리포트 생성: {output_path}"
    
    def _create_summary_sheet(self, ws):
        """요약 시트 생성"""
        try:
            from openpyxl.styles import Font, PatternFill, Alignment
        except:
            return
        
        # 제목
        ws['A1'] = "GAIM Lab v3.0 분석 보고서"
        ws['A1'].font = Font(size=16, bold=True, color="FFFFFF")
        ws['A1'].fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
        ws.merge_cells('A1:D1')
        
        # 생성 일시
        ws['A2'] = f"생성일: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        ws['A2'].font = Font(size=10, italic=True)
        
        # 통계
        avg_score = sum(r["total_score"] for r in self.results) / len(self.results)
        max_score = max(r["total_score"] for r in self.results)
        min_score = min(r["total_score"] for r in self.results)
        
        row = 4
        ws[f'A{row}'] = "통계 요약"
        ws[f'A{row}'].font = Font(bold=True, size=12, color="FFFFFF")
        ws[f'A{row}'].fill = PatternFill(start_color="70AD47", end_color="70AD47", fill_type="solid")
        ws.merge_cells(f'A{row}:B{row}')
        
        row = 5
        stats = [
            ("분석 영상 수", len(self.results)),
            ("평균 점수", f"{avg_score:.1f}"),
            ("최고 점수", f"{max_score:.1f}"),
            ("최저 점수", f"{min_score:.1f}"),
        ]
        
        for label, value in stats:
            ws[f'A{row}'] = label
            ws[f'B{row}'] = value
            ws[f'A{row}'].font = Font(bold=True)
            row += 1
        
        # 등급 분포
        row += 1
        ws[f'A{row}'] = "등급 분포"
        ws[f'A{row}'].font = Font(bold=True, size=12, color="FFFFFF")
        ws[f'A{row}'].fill = PatternFill(start_color="70AD47", end_color="70AD47", fill_type="solid")
        ws.merge_cells(f'A{row}:B{row}')
        
        grades = {}
        for r in self.results:
            g = r["grade_letter"]
            grades[g] = grades.get(g, 0) + 1
        
        row = row + 1
        for grade in sorted(grades.keys(), reverse=True):
            ws[f'A{row}'] = f"{grade} 등급"
            ws[f'B{row}'] = grades[grade]
            row += 1
        
        # 열 너비 설정
        ws.column_dimensions['A'].width = 20
        ws.column_dimensions['B'].width = 15
    
    def _create_detail_sheet(self, ws):
        """상세 분석 시트"""
        try:
            from openpyxl.styles import Font, PatternFill
        except:
            return
        
        # 헤더
        headers = ["순번", "영상명", "과목", "교사", "학년", "점수", "등급"]
        for col, header in enumerate(headers, 1):
            cell = ws.cell(1, col)
            cell.value = header
            cell.font = Font(bold=True, color="FFFFFF")
            cell.fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
        
        # 데이터
        sorted_results = sorted(self.results, key=lambda x: x["total_score"], reverse=True)
        for row_idx, result in enumerate(sorted_results, 2):
            ws.cell(row_idx, 1).value = row_idx - 1
            ws.cell(row_idx, 2).value = result["video"]
            ws.cell(row_idx, 3).value = result["subject"]
            ws.cell(row_idx, 4).value = result["teacher"]
            ws.cell(row_idx, 5).value = result["grade"]
            ws.cell(row_idx, 6).value = result["total_score"]
            ws.cell(row_idx, 7).value = result["grade_letter"]
        
        # 열 너비
        ws.column_dimensions['A'].width = 8
        ws.column_dimensions['B'].width = 35
        ws.column_dimensions['C'].width = 10
        ws.column_dimensions['D'].width = 10
        ws.column_dimensions['E'].width = 8
        ws.column_dimensions['F'].width = 10
        ws.column_dimensions['G'].width = 8
    
    def _create_teacher_sheet(self, ws):
        """교사별 분석 시트"""
        try:
            from openpyxl.styles import Font, PatternFill
        except:
            return
        
        # 교사별 평균
        teacher_stats = {}
        for result in self.results:
            teacher = result["teacher"]
            if teacher not in teacher_stats:
                teacher_stats[teacher] = {"scores": [], "count": 0}
            teacher_stats[teacher]["scores"].append(result["total_score"])
            teacher_stats[teacher]["count"] += 1
        
        # 헤더
        headers = ["교사명", "수업 개수", "평균 점수", "최고 점수", "최저 점수"]
        for col, header in enumerate(headers, 1):
            cell = ws.cell(1, col)
            cell.value = header
            cell.font = Font(bold=True, color="FFFFFF")
            cell.fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
        
        # 데이터
        for row_idx, (teacher, stats) in enumerate(
            sorted(teacher_stats.items(), key=lambda x: sum(x[1]["scores"]) / len(x[1]["scores"]), reverse=True),
            2
        ):
            ws.cell(row_idx, 1).value = teacher
            ws.cell(row_idx, 2).value = stats["count"]
            ws.cell(row_idx, 3).value = round(sum(stats["scores"]) / len(stats["scores"]), 1)
            ws.cell(row_idx, 4).value = max(stats["scores"])
            ws.cell(row_idx, 5).value = min(stats["scores"])
        
        # 열 너비
        for col in range(1, 6):
            ws.column_dimensions[chr(64 + col)].width = 15
    
    def _create_subject_sheet(self, ws):
        """과목별 분석 시트"""
        try:
            from openpyxl.styles import Font, PatternFill
        except:
            return
        
        # 과목별 통계
        subject_stats = {}
        for result in self.results:
            subject = result["subject"]
            if subject not in subject_stats:
                subject_stats[subject] = {"scores": [], "count": 0}
            subject_stats[subject]["scores"].append(result["total_score"])
            subject_stats[subject]["count"] += 1
        
        # 헤더
        headers = ["과목", "영상 수", "평균 점수", "최고 점수", "최저 점수"]
        for col, header in enumerate(headers, 1):
            cell = ws.cell(1, col)
            cell.value = header
            cell.font = Font(bold=True, color="FFFFFF")
            cell.fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
        
        # 데이터
        for row_idx, (subject, stats) in enumerate(
            sorted(subject_stats.items(), key=lambda x: sum(x[1]["scores"]) / len(x[1]["scores"]), reverse=True),
            2
        ):
            ws.cell(row_idx, 1).value = subject
            ws.cell(row_idx, 2).value = stats["count"]
            ws.cell(row_idx, 3).value = round(sum(stats["scores"]) / len(stats["scores"]), 1)
            ws.cell(row_idx, 4).value = max(stats["scores"])
            ws.cell(row_idx, 5).value = min(stats["scores"])
        
        # 열 너비
        for col in range(1, 6):
            ws.column_dimensions[chr(64 + col)].width = 15
    
    def generate_pdf_report(self, filename: str = "analysis_report.pdf") -> str:
        """
        PDF 리포트 생성
        
        Args:
            filename: 출력 파일명
            
        Returns:
            생성된 파일 경로
        """
        try:
            from reportlab.lib.pagesizes import A4, letter
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import inch
            from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, PageBreak
            from reportlab.lib import colors
        except ImportError:
            return "❌ reportlab 패키지가 필요합니다: pip install reportlab"
        
        output_path = self.output_dir / filename
        
        # PDF 생성
        doc = SimpleDocTemplate(str(output_path), pagesize=letter)
        story = []
        styles = getSampleStyleSheet()
        
        # 제목
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            textColor=colors.HexColor('#4472C4'),
            spaceAfter=12,
            alignment=1  # 중앙 정렬
        )
        story.append(Paragraph("GAIM Lab v3.0<br/>교육 수업 분석 보고서", title_style))
        story.append(Spacer(1, 0.3*inch))
        
        # 생성 일시
        story.append(Paragraph(f"생성일: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}", styles['Normal']))
        story.append(Spacer(1, 0.2*inch))
        
        # 통계
        avg_score = sum(r["total_score"] for r in self.results) / len(self.results)
        max_score = max(r["total_score"] for r in self.results)
        min_score = min(r["total_score"] for r in self.results)
        
        stats_data = [
            ["분석 영상 수", f"{len(self.results)}개"],
            ["평균 점수", f"{avg_score:.1f}점"],
            ["최고 점수", f"{max_score:.1f}점"],
            ["최저 점수", f"{min_score:.1f}점"],
        ]
        
        stats_table = Table(stats_data)
        stats_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (1, 0), colors.HexColor('#4472C4')),
            ('TEXTCOLOR', (0, 0), (1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (1, 0), 12),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        story.append(stats_table)
        story.append(Spacer(1, 0.3*inch))
        
        # 상세 표
        story.append(Paragraph("분석 결과 상세", styles['Heading2']))
        story.append(Spacer(1, 0.1*inch))
        
        sorted_results = sorted(self.results, key=lambda x: x["total_score"], reverse=True)
        detail_data = [["순번", "영상명", "과목", "점수", "등급"]]
        
        for idx, result in enumerate(sorted_results[:10], 1):  # 상위 10개만
            detail_data.append([
                str(idx),
                result["video"][:30],
                result["subject"],
                f"{result['total_score']:.1f}",
                result["grade_letter"]
            ])
        
        detail_table = Table(detail_data)
        detail_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#4472C4')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 10),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('GRID', (0, 0), (-1, -1), 1, colors.black),
            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#F0F0F0')])
        ]))
        story.append(detail_table)
        
        # PDF 생성
        doc.build(story)
        
        return f"✅ PDF 리포트 생성: {output_path}"
    
    def generate_powerpoint_report(self, filename: str = "analysis_report.pptx") -> str:
        """
        PowerPoint 리포트 생성
        
        Args:
            filename: 출력 파일명
            
        Returns:
            생성된 파일 경로
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor
        except ImportError:
            return "❌ python-pptx 패키지가 필요합니다: pip install python-pptx"
        
        # 프레젠테이션 생성
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # 슬라이드 1: 제목
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(8)
        height = Inches(2)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = "GAIM Lab v3.0"
        title_frame.paragraphs[0].font.size = Pt(54)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(68, 114, 196)
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        subtitle_box = slide.shapes.add_textbox(left, top + Inches(1.5), width, Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "교육 수업 분석 보고서"
        subtitle_frame.paragraphs[0].font.size = Pt(32)
        subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(112, 173, 71)
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 슬라이드 2: 통계
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "분석 통계"
        
        avg_score = sum(r["total_score"] for r in self.results) / len(self.results)
        max_score = max(r["total_score"] for r in self.results)
        min_score = min(r["total_score"] for r in self.results)
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.text = f"분석 영상: {len(self.results)}개"
        
        p = text_frame.add_paragraph()
        p.text = f"평균 점수: {avg_score:.1f}점"
        p.level = 0
        
        p = text_frame.add_paragraph()
        p.text = f"최고 점수: {max_score:.1f}점"
        p.level = 0
        
        p = text_frame.add_paragraph()
        p.text = f"최저 점수: {min_score:.1f}점"
        p.level = 0
        
        # 슬라이드 3: 상위 5개
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "상위 5개 영상"
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        sorted_results = sorted(self.results, key=lambda x: x["total_score"], reverse=True)
        for idx, result in enumerate(sorted_results[:5], 1):
            p = text_frame.add_paragraph()
            p.text = f"{idx}. {result['video']} - {result['total_score']:.1f}점 ({result['grade_letter']})"
            p.level = 0
        
        # 파일 저장
        output_path = self.output_dir / filename
        prs.save(output_path)
        
        return f"✅ PowerPoint 리포트 생성: {output_path}"


if __name__ == "__main__":
    # 테스트
    from pathlib import Path
    import json
    
    results_file = Path(__file__).parent.parent / "output" / "batch_analysis_20260207_005428" / "batch_results.json"
    output_dir = Path(__file__).parent.parent / "output" / "batch_analysis_20260207_005428"
    
    if results_file.exists():
        with open(results_file, encoding="utf-8") as f:
            results = json.load(f)
        
        generator = ReportGenerator(results, str(output_dir))
        
        print("=" * 70)
        print("📊 다중 형식 리포트 생성")
        print("=" * 70)
        
        print("\n" + generator.generate_excel_report("advanced_report.xlsx"))
        print(generator.generate_pdf_report("advanced_report.pdf"))
        print(generator.generate_powerpoint_report("advanced_report.pptx"))
    else:
        print("❌ batch_results.json 파일을 찾을 수 없습니다.")
